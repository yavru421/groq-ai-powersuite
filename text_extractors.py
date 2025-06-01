import PyPDF2
import docx
from pptx import Presentation
import csv
import json
import os
from PIL import Image
import io
import base64

def extract_text_from_file(uploaded_file_path: str) -> str:
    """
    Extracts text from various file types.
    :param uploaded_file_path: Path to the uploaded file.
    :return: Extracted text as a string, or an error message.
    """
    _, file_extension = os.path.splitext(uploaded_file_path)
    file_extension = file_extension.lower()

    try:
        if file_extension == '.pdf':
            with open(uploaded_file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                return text
        elif file_extension == '.docx':
            doc = docx.Document(uploaded_file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text
        elif file_extension == '.pptx':
            prs = Presentation(uploaded_file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Safely extract text from shapes that have a text attribute
                    shape_text = getattr(shape, "text", None)
                    if shape_text:
                        text += shape_text + "\n"
            return text
        elif file_extension == '.txt':
            with open(uploaded_file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif file_extension == '.csv':
            with open(uploaded_file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                return "\n".join([",".join(row) for row in reader])
        elif file_extension == '.json':
            with open(uploaded_file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        else:
            return "Unsupported file type or error extracting text."
    except Exception as e:
        return f"Error extracting text from {uploaded_file_path}: {str(e)}"

def extract_image_bytes_and_base64(uploaded_file_path: str, make_square=False, target_size_kb=768):
    """
    Reads an image file, converts it to JPEG, resizes if necessary to be under a target size,
    and returns its base64 representation and original bytes.
    """
    try:
        img = Image.open(uploaded_file_path)
        if img.mode == 'RGBA' or img.mode == 'P':
            img = img.convert('RGB')

        if make_square:
            size = (max(img.size), max(img.size))
            new_img = Image.new("RGB", size, (255, 255, 255))
            new_img.paste(img, (int((size[0] - img.size[0]) / 2), int((size[1] - img.size[1]) / 2)))
            img = new_img
            
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='JPEG', quality=90)
        
        quality = 85
        while img_byte_arr.tell() / 1024 > target_size_kb and quality > 10:
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='JPEG', quality=quality)
            quality -= 5
            
        original_bytes = img_byte_arr.getvalue()
        base64_image = base64.b64encode(original_bytes).decode('utf-8')
        
        return original_bytes, f"data:image/jpeg;base64,{base64_image}"
    except Exception as e:
        print(f"Error processing image: {e}")
        return None, None

# Wrapper for Gradio File object
def extract_text_from_gradio_file(gradio_file_obj) -> str:
    if gradio_file_obj is None:
        return "No file provided."
    return extract_text_from_file(gradio_file_obj.name) # .name gives the temp file path

def extract_image_data_from_gradio_image(gradio_image_obj_path: str):
    if gradio_image_obj_path is None:
        return None, None, "No image provided."
    
    # Gradio gr.Image with type="filepath" returns the path to a temporary file
    original_bytes, b64_string = extract_image_bytes_and_base64(gradio_image_obj_path)
    
    if b64_string:
        return original_bytes, b64_string, "Image processed successfully."
    else:
        return None, None, "Error processing image."
